import argparse
import ast
from difflib import SequenceMatcher
import html
import json
import os
import re
import smtplib
import ssl
import urllib.parse
import urllib.request
import xml.etree.ElementTree as ET
from datetime import datetime
from email.message import EmailMessage
from html.parser import HTMLParser
from pathlib import Path


BASE_DIR = Path(__file__).resolve().parent
SETTINGS_PATH = BASE_DIR / "settings.json"
DEFAULT_FEED_URL = "https://news.google.com/rss?hl=ko&gl=KR&ceid=KR:ko"
LAST_GEMINI_ERROR = ""


class NewsItem:
    def __init__(
        self,
        title: str,
        link: str,
        source: str,
        summary: str = "",
        article_text: str = "",
    ) -> None:
        self.title = title
        self.link = link
        self.source = source
        self.summary = summary
        self.article_text = article_text


class ContentPackage:
    def __init__(
        self,
        news_title: str,
        blog_post: str,
        tistory_post: str,
        thread_post: str,
        slide_script: str,
        vrew_script: str,
        directory: str = "",
        pptx_path: str = "",
    ) -> None:
        self.news_title = news_title
        self.blog_post = blog_post
        self.tistory_post = tistory_post
        self.thread_post = thread_post
        self.slide_script = slide_script
        self.vrew_script = vrew_script
        self.directory = directory
        self.pptx_path = pptx_path


def load_env_file(path: Path = BASE_DIR / ".env") -> None:
    if not path.exists():
        return

    for raw_line in path.read_text(encoding="utf-8-sig").splitlines():
        line = raw_line.strip()
        if not line or line.startswith("#") or "=" not in line:
            continue
        key, value = line.split("=", 1)
        os.environ.setdefault(key.strip(), value.strip().strip('"').strip("'"))


def load_settings(path: Path = SETTINGS_PATH) -> dict:
    if not path.exists():
        return {}
    return json.loads(path.read_text(encoding="utf-8-sig"))


def as_bool(value, default: bool = False) -> bool:
    if isinstance(value, bool):
        return value
    if value is None:
        return default
    return str(value).strip().lower() in {"1", "true", "yes", "y", "on"}


def get_secret(name: str) -> str:
    value = os.getenv(name, "").strip()
    if value:
        return value

    try:
        import streamlit as st

        return str(st.secrets.get(name, "")).strip()
    except Exception:
        return ""


class ArticleTextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.skip_depth = 0
        self.parts: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in {"script", "style", "noscript", "svg", "header", "footer", "nav"}:
            self.skip_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style", "noscript", "svg", "header", "footer", "nav"} and self.skip_depth:
            self.skip_depth -= 1

    def handle_data(self, data: str) -> None:
        text = " ".join(data.split())
        if self.skip_depth or len(text) < 35:
            return
        self.parts.append(text)


class LinkParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.links: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag != "a":
            return
        for name, value in attrs:
            if name == "href" and value:
                self.links.append(html.unescape(value))


def strip_html(value: str) -> str:
    parser = ArticleTextParser()
    parser.feed(value or "")
    return " ".join(parser.parts).strip()


def extract_links(value: str) -> list[str]:
    parser = LinkParser()
    parser.feed(value or "")
    return parser.links


def is_google_news_url(url: str) -> bool:
    host = urllib.parse.urlparse(url).netloc.lower()
    return "news.google." in host or host == "google.com" or host.endswith(".google.com")


def is_useful_article_text(text: str, title: str) -> bool:
    normalized_text = " ".join((text or "").split())
    normalized_title = " ".join((title or "").split())
    if len(normalized_text) < 180:
        return False
    if normalized_text == normalized_title:
        return False
    if normalized_title and normalized_text.startswith(normalized_title) and len(normalized_text) < len(normalized_title) + 80:
        return False
    return True


def article_score(item: NewsItem) -> int:
    score = 0
    if is_useful_article_text(item.article_text, item.title):
        score += min(len(item.article_text), 3000)
    if is_useful_article_text(item.summary, item.title):
        score += min(len(item.summary), 800)
    if item.article_text:
        score += 200
    return score


def select_content_item(items: list[NewsItem], preferred_index: int = 0) -> tuple[NewsItem, int, bool]:
    if not items:
        raise ValueError("?좏깮???댁뒪媛 ?놁뒿?덈떎.")

    if preferred_index > 0:
        index = min(max(preferred_index, 1), len(items))
        return items[index - 1], index, False

    scored_items = [(article_score(item), index, item) for index, item in enumerate(items, start=1)]
    scored_items.sort(key=lambda value: value[0], reverse=True)
    best_score, best_index, best_item = scored_items[0]
    if best_score <= 0:
        return items[0], 1, True
    return best_item, best_index, True


def extract_json_object(value: str) -> dict:
    text = (value or "").strip()
    if text.startswith("```"):
        text = re.sub(r"^```(?:json)?", "", text, flags=re.IGNORECASE).strip()
        text = re.sub(r"```$", "", text).strip()

    def remove_bad_control_chars(raw: str) -> str:
        return "".join(
            char
            for char in raw
            if char in "\n\r\t" or ord(char) >= 32
        )

    def escape_string_newlines(raw: str) -> str:
        result: list[str] = []
        in_string = False
        escaped = False
        for char in raw:
            if escaped:
                result.append(char)
                escaped = False
                continue
            if char == "\\":
                result.append(char)
                escaped = True
                continue
            if char == '"':
                in_string = not in_string
                result.append(char)
                continue
            if in_string and char in "\n\r":
                result.append("\\n")
                continue
            if in_string and char == "\t":
                result.append("\\t")
                continue
            result.append(char)
        return "".join(result)

    def loads_lenient(raw: str) -> dict:
        cleaned = remove_bad_control_chars(raw).strip()
        try:
            return json.loads(cleaned)
        except json.JSONDecodeError:
            return json.loads(escape_string_newlines(cleaned))

    try:
        return loads_lenient(text)
    except json.JSONDecodeError:
        match = re.search(r"\{.*\}", text, flags=re.DOTALL)
        if not match:
            raise
        return loads_lenient(match.group(0))


def normalize_generated_text_field(value) -> str:
    if isinstance(value, dict):
        title = str(value.get("title") or value.get("headline") or "").strip()
        content = str(
            value.get("content")
            or value.get("body")
            or value.get("text")
            or value.get("post")
            or ""
        ).strip()
        if title and content and not content.lstrip().startswith("#"):
            return f"# {title}\n\n{content}".strip()
        return (content or title).strip()
    if isinstance(value, list):
        return "\n\n".join(normalize_generated_text_field(item) for item in value if item).strip()

    text = str(value or "").strip()
    if text.startswith("{") and ("'content'" in text or '"content"' in text or "'body'" in text or '"body"' in text):
        try:
            parsed = ast.literal_eval(text)
            return normalize_generated_text_field(parsed)
        except (ValueError, SyntaxError):
            pass
    return text.replace("\\r\\n", "\n").replace("\\n", "\n").replace("\\t", "\t").strip()


def format_tistory_post(text: str) -> str:
    text = (text or "").strip()
    if not text:
        return ""

    lines = [line.strip() for line in text.splitlines()]
    formatted: list[str] = []
    for line in lines:
        if not line:
            if formatted and formatted[-1] != "":
                formatted.append("")
            continue

        is_heading = line.startswith("#") or re.match(r"^\d+\.\s+", line)
        is_list = line.startswith(("- ", "* ", "[ ]", "- [ ]"))

        if is_heading and formatted and formatted[-1] != "":
            formatted.append("")
        formatted.append(line)
        if is_heading:
            formatted.append("")

    text = "\n".join(formatted)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()

    if "\n\n" not in text and len(text) > 700:
        sentences = re.split(r"(?<=[.!?。！？])\s+", text)
        chunks: list[str] = []
        for index in range(0, len(sentences), 3):
            chunk = " ".join(sentence.strip() for sentence in sentences[index:index + 3] if sentence.strip())
            if chunk:
                chunks.append(chunk)
        text = "\n\n".join(chunks)

    return text


def create_grounded_article_context(item: NewsItem, article_context: str) -> str:
    api_key = get_secret("GEMINI_API_KEY")
    if not api_key:
        return ""

    try:
        from google import genai
        from google.genai import types
    except Exception:
        return ""

    model = get_secret("GEMINI_GROUNDING_MODEL") or get_secret("GEMINI_MODEL") or "gemini-3.5-flash"
    prompt = f"""
Write in Korean.
Use Google Search to strengthen the news context below.
Do not invent exact numbers, quotes, dates, company names, or policy names unless they are found in sources.
If the original article text is limited, use search results to explain confirmed facts, background, likely impact, and what readers should verify next.
Return only a concise briefing, not JSON.

News title:
{item.title}

Publisher/source:
{item.source}

Original link:
{item.link}

Collected text:
{article_context}

Required structure:
1. Confirmed facts
2. Background
3. Why it matters
4. What to watch next
5. Source notes
"""

    try:
        client = genai.Client(api_key=api_key)
        grounding_tool = types.Tool(google_search=types.GoogleSearch())
        response = client.models.generate_content(
            model=model,
            contents=prompt,
            config=types.GenerateContentConfig(tools=[grounding_tool]),
        )
    except Exception:
        return ""

    text = " ".join((response.text or "").split())
    if len(text) < 180:
        return ""
    return text[:4000]


def create_gemini_content(item: NewsItem, article_context: str, use_grounding: bool = False) -> dict | None:
    global LAST_GEMINI_ERROR
    LAST_GEMINI_ERROR = ""
    api_key = get_secret("GEMINI_API_KEY")
    if not api_key:
        LAST_GEMINI_ERROR = "GEMINI_API_KEY가 없습니다."
        return None

    try:
        from google import genai
        from google.genai import types
    except Exception as exc:
        LAST_GEMINI_ERROR = f"Gemini 패키지를 불러오지 못했습니다: {exc}"
        return None

    model = get_secret("GEMINI_MODEL") or "gemini-2.5-pro"
    prompt = f"""
You are a Korean blog columnist for 50s and 60s readers.
Write the final content in Korean.
Do not write a report, checklist-only memo, or plain news summary. Turn the news topic into a practical Korean blog column.

Most important rules for blog_post:
- Write a complete Korean blog post of 2800 to 3500 Korean characters.
- Use the style of a practical 5060 money/life strategy blog.
- Start with a strong title like "[?듭떖 ?ㅼ썙?? 吏湲?以鍮꾪빐?????댁쑀".
- Opening should directly tell readers why this issue matters to their money, retirement, work, household budget, or future choices.
- Use numbered sections, for example "1. ?댁뒪 ?듭떖", "2. ??以묒슂?쒓?", "3. ?쒖옣怨??앺솢??誘몄튌 ?곹뼢", "4. ?욎쑝濡?蹂??ъ씤??, "5. 泥댄겕由ъ뒪??.
- Do not include fictional case studies or hypothetical personal examples unless the article itself provides them.
- Include a checklist near the end with 4 to 6 action items.
- End with a conclusion and a short "?ㅼ쓬???뺤씤??寃? preview.
- If article text is limited, do not mention that limitation to readers. Use the title, source, verified context, and general background to write a useful blog-style interpretation.
- Never use these phrases in blog_post: "???? ??? ??", "??? ??", "?? ??", "?? ??", "??? ??", "???? ???", "?? ??".
- Do not invent exact numbers, quotes, schedules, company names, or policy names that are not provided.
- The writing should be confident, practical, and easy to read. Avoid vague filler such as "吏耳쒕킄???⑸땲?? repeated too often.

Return only one valid JSON object. Use exactly these keys:
blog_post, tistory_post, thread_post, slide_script, vrew_script

Requirements:
- blog_post: 2800 to 3500 Korean characters, news commentary blog article with numbered sections and checklist.
- tistory_post: adapt blog_post for Tistory with SEO-friendly title, H2/H3 headings, conclusion, and 5 to 8 tags.
  Format it for easy reading on Tistory:
  use Markdown headings, put one blank line after every heading,
  split long text into short paragraphs of 2 to 4 sentences,
  put one blank line between all paragraphs,
  use bullet/checklist lines on separate lines,
  never write sections as one long connected block.
- thread_post: 250 to 330 Korean characters, like a first social media post.
- slide_script: 6 YouTube slides. Use "## ???? 1." through "## ???? 6." in order. Each slide must include "- ?? ??:" and "- ????:".
- vrew_script: 6 scene narration script. Use "## ???? 1." through "## ???? 6." in order.

News title:
{item.title}

Source:
{item.source}

Original link:
{item.link}

Collected article context:
{article_context}
"""

    client = genai.Client(api_key=api_key)

    def request_content(with_grounding: bool):
        if with_grounding:
            grounding_tool = types.Tool(google_search=types.GoogleSearch())
            return client.models.generate_content(
                model=model,
                contents=(
                    prompt
                    + "\n\nAdditional instruction: Use Google Search grounding to verify and enrich the context before writing. "
                    "Then write it as a practical news commentary blog column with numbered sections, a reader checklist, and a closing preview. "
                    "Keep the final JSON in Korean."
                ),
                config=types.GenerateContentConfig(
                    tools=[grounding_tool],
                    temperature=0.85,
                    max_output_tokens=8192,
                ),
            )
        return client.models.generate_content(
            model=model,
            contents=prompt,
            config=types.GenerateContentConfig(
                temperature=0.85,
                max_output_tokens=8192,
            ),
        )

    try:
        try:
            response = request_content(use_grounding)
        except Exception as grounding_exc:
            if not use_grounding:
                raise
            LAST_GEMINI_ERROR = f"Google Search grounding ?ㅽ뙣, ?쇰컲 Gemini濡??ъ떆?꾪뻽?듬땲?? {grounding_exc}"
            response = request_content(False)
        data = extract_json_object(response.text)
    except Exception as exc:
        LAST_GEMINI_ERROR = f"Gemini ?앹꽦 ?ㅽ뙣: {exc}"
        return None

    blog_post = str(data.get("blog_post", "")).strip()[:4200]
    tistory_post = format_tistory_post(str(data.get("tistory_post", "")).strip())[:4500]
    thread_post = ensure_thread_length(str(data.get("thread_post", "")).strip(), item.title, article_context)
    slide_script = str(data.get("slide_script", "")).strip()
    vrew_script = str(data.get("vrew_script", "")).strip()

    if len(blog_post) < 2500:
        LAST_GEMINI_ERROR = f"Gemini blog post was too short: {len(blog_post)} chars"
        return None

    forbidden_blog_phrases = [
        "기사에서 확인한 내용",
        "정보가 부족",
        "원문 확인",
        "자동 수집",
        "확인한 범위",
        "보강해야 합니다",
        "구체적인 수치",
        "기사 본문",
    ]
    if any(phrase in blog_post for phrase in forbidden_blog_phrases):
        LAST_GEMINI_ERROR = "Gemini媛 釉붾줈洹?湲 ???蹂닿퀬?쒖떇 臾몄옣??留뚮뱾?덉뒿?덈떎."
        return None

    if not all([blog_post, tistory_post, thread_post, slide_script, vrew_script]):
        LAST_GEMINI_ERROR = "Gemini ?묐떟?먯꽌 ?꾩슂????ぉ ?쇰?媛 鍮꾩뼱 ?덉뒿?덈떎."
        return None

    return {
        "blog_post": blog_post,
        "tistory_post": tistory_post,
        "thread_post": thread_post,
        "slide_script": slide_script,
        "vrew_script": vrew_script,
    }


def create_derivative_content_from_blog(
    blog_post: str,
    title: str = "",
    source: str = "직접 작성",
    link: str = "",
) -> dict[str, str]:
    title = title.strip() or extract_title_from_text(blog_post, "직접 작성한 블로그 글")
    source = source.strip() or "직접 작성"
    link = link.strip()
    api_key = get_secret("GEMINI_API_KEY")
    data: dict[str, str] | None = None

    if api_key:
        try:
            from google import genai

            model = get_secret("GEMINI_MODEL") or "gemini-2.5-pro"
            prompt = f"""
?덈뒗 ?쒓뎅??釉붾줈洹??륂뤌 肄섑뀗痢??묎???
?꾨옒 釉붾줈洹?湲??湲곗? ?먭퀬濡??쇱븘 ?곗뒪?좊━ 湲, ?곕젅?? ?좏뒠釉??щ씪?대뱶 ?蹂? Vrew ?蹂몄쓣 ?ㅼ떆 ?묒꽦?섎씪.
蹂닿퀬?쒖쿂???붿빟?섏? 留먭퀬, ?낆옄媛 蹂닿퀬 ?띠뼱吏???쒕ぉ怨??먯뿰?ㅻ윭??留먰닾濡?媛곸깋?섎씪.
釉붾줈洹?湲???녿뒗 援ъ껜???レ옄, 諛쒖뼵, ?쇱젙, 湲곗뾽紐낆? 吏?대궡吏 留덈씪.
異쒕젰? JSON 媛앹껜 ?섎굹留?諛섑솚?섎씪.
?ㅻ뒗 tistory_post, thread_post, slide_script, vrew_script ??媛쒕쭔 ?ъ슜?섎씪.

議곌굔:
- tistory_post: ?곗뒪?좊━ 釉붾줈洹몄슜 媛곸깋 湲. 寃???좎엯???쒕ぉ, ?먯뿰?ㅻ윭???꾩엯, H2/H3 ?뚯젣紐? 紐⑸줉, 留덈Т由? 愿???쒓렇 5~8媛??ы븿.
  ?곗뒪?좊━?먯꽌 ?쎄린 醫뗪쾶 ?쒕ぉ怨??뚯젣紐??ㅼ뿉??諛섎뱶??鍮?以꾩쓣 ?ｊ퀬, 湲?臾몄옣? 2~4臾몄옣 ?⑥쐞??吏㏃? 臾몃떒?쇰줈 ?섎늻?대씪.
  紐⑤뱺 臾몃떒 ?ъ씠?먮뒗 鍮?以꾩쓣 ?섎굹 ?ｌ뼱??
  泥댄겕由ъ뒪?몄? 紐⑸줉? ??以꾩뿉 ?섎굹???곌퀬, 湲 ?꾩껜瑜??덈? ???⑹뼱由щ줈 ?댁뼱 ?곗? 留덈씪.
- thread_post: 280~330?? SNS 泥?湲泥섎읆 沅곴툑利앹쓣 留뚮뱾 寃?
- slide_script: ?좏뒠釉??쒖옉???щ씪?대뱶 6??援ъ꽦. 諛섎뱶??"## ?щ씪?대뱶 1."遺??"## ?щ씪?대뱶 6."源뚯? ?꾩뿉???꾨옒濡??쒖꽌?濡??묒꽦?섍퀬, 媛??λ쭏??"- ?붾㈃ 臾멸뎄:"? "- ?대젅?댁뀡:" ?ы븿. ?붾㈃ 臾멸뎄??吏㏐퀬 媛뺥븯寃?
- vrew_script: Vrew??遺숈뿬?ｊ린 醫뗭? ?먯뿰?ㅻ윭??留먰닾???λ㈃蹂??대젅?댁뀡 ?蹂? 諛섎뱶??"## ?щ씪?대뱶 1."遺??"## ?щ씪?대뱶 6."源뚯? ?꾩뿉???꾨옒濡??쒖꽌?濡??묒꽦

?쒕ぉ:
{title}

?먮Ц 異쒖쿂:
{source}

?먮Ц 留곹겕:
{link}

湲곗? 釉붾줈洹?湲:
{blog_post[:6000]}
"""
            client = genai.Client(api_key=api_key)
            response = client.models.generate_content(model=model, contents=prompt)
            raw_data = extract_json_object(response.text)
            data = {
                "tistory_post": format_tistory_post(str(raw_data.get("tistory_post", "")).strip())[:4500],
                "thread_post": ensure_thread_length(str(raw_data.get("thread_post", "")).strip(), title, blog_post),
                "slide_script": str(raw_data.get("slide_script", "")).strip(),
                "vrew_script": str(raw_data.get("vrew_script", "")).strip(),
            }
        except Exception:
            data = None

    item = NewsItem(title=title, link=link, source=source, summary=blog_post[:500], article_text=blog_post)
    if not data or not all(data.values()):
        data = fallback_derivative_content_from_blog(blog_post, item)

    slide_blocks = ensure_six_slide_blocks(data["slide_script"], item, blog_post)
    data["slide_script"] = format_slide_script(slide_blocks, source, link)
    data["vrew_script"] = format_vrew_script(slide_blocks, source)
    data["thread_post"] = ensure_thread_length(data["thread_post"], title, blog_post)
    data["tistory_post"] = format_tistory_post(data["tistory_post"])[:4500]
    return data


def ensure_thread_length(thread_post: str, title: str, context: str) -> str:
    thread_post = " ".join((thread_post or "").split())
    if len(thread_post) >= 240:
        return thread_post[:330]

    context_preview = " ".join((context or "").split())[:130]
    expanded = (
        f"{thread_post} "
        f"?듭떖? '{title}' ?댁뒋瑜??⑥닚???댁뒪濡??섍린吏 ?딄퀬, 諛곌꼍怨??ㅼ젣 ?곹뼢, ?ㅼ쓬 ?吏곸엫源뚯? ?④퍡 蹂대뒗 寃껋엯?덈떎. "
        f"{context_preview} "
        "吏湲덉? 寃곕줎蹂대떎 ?뺤씤???ъ씤?몃? ?뺣━?대몢??寃껋씠 以묒슂?⑸땲??"
    )
    return " ".join(expanded.split())[:330]


def extract_title_from_text(text: str, fallback: str) -> str:
    for line in text.splitlines():
        cleaned = line.strip()
        if cleaned.startswith("#"):
            return cleaned.lstrip("# ").strip()[:90] or fallback
    first_line = next((line.strip() for line in text.splitlines() if line.strip()), "")
    return first_line[:90] or fallback


def fallback_derivative_content_from_blog(blog_post: str, item: NewsItem) -> dict[str, str]:
    title = item.title
    preview = " ".join(blog_post.split())[:900]
    tistory_post = f"""# {title}

## ?듭떖 ?붿빟

{preview}

## ??二쇰ぉ?댁빞 ?좉퉴?

??湲???듭떖? ?⑥닚???댁뒪 ?꾨떖???꾨땲?? ?낆옄媛 吏湲??뺤씤?댁빞 ??蹂?붿쓽 ?먮쫫???뺣━?섎뒗 ???덉뒿?덈떎. 蹂몃Ц?먯꽌 ?ㅻ， 諛곌꼍怨??곹뼢, ?욎쑝濡쒖쓽 ?吏곸엫???섎늻??蹂대㈃ ??紐낇솗?섍쾶 ?댄빐?????덉뒿?덈떎.

## ?뺤씤???ъ씤??

- ???댁뒋媛 ?쒖옉??諛곌꼍
- 愿???쒖옣怨??앺솢??誘몄튌 ?곹뼢
- ?욎쑝濡??섏삱 ?꾩냽 諛쒗몴??諛섏쓳

## 留덈Т由?

?뺥솗???먮떒???꾪빐?쒕뒗 ?먮Ц怨?異붽? ?먮즺瑜??④퍡 ?뺤씤?섎뒗 寃껋씠 醫뗭뒿?덈떎. 吏湲덉? 蹂?붿쓽 諛⑺뼢??癒쇱? ?댄빐?섍퀬, ?ㅼ쓬 ?吏곸엫??李⑤텇??吏耳쒕낵 ?쒖젏?낅땲??

### 愿???쒓렇

#?댁뒪?뺣━ #?댁뒋遺꾩꽍 #寃쎌젣?댁뒪 #肄섑뀗痢좎젣??#?곗뒪?좊━ #釉붾줈洹멸??곌린
"""
    thread_post = (
        f"{title} ???댁뒋???⑥닚???댁뒪 ??以꾨낫???욎쑝濡쒖쓽 ?먮쫫??蹂댁뿬二쇰뒗 ?좏샇??媛源앹뒿?덈떎. "
        "?듭떖? 諛곌꼍, ?ㅼ젣 ?곹뼢, ?ㅼ쓬 ?吏곸엫???④퍡 蹂대뒗 寃껋엯?덈떎. 吏湲??뺤씤???ъ씤?몃? ?뺣━?대몢硫??댄썑 蹂?붿뿉 ??鍮좊Ⅴ寃???묓븷 ???덉뒿?덈떎."
    )
    slide_script = f"""# ?좏뒠釉??щ씪?대뱶 ?蹂?
## ?щ씪?대뱶 1. ?ㅽ봽??
- ?붾㈃ 臾멸뎄: 吏湲??뺤씤?댁빞 ???댁뒋
- ?대젅?댁뀡: ?ㅻ뒛? "{title}" ?댁뒋瑜??듭떖留?鍮좊Ⅴ寃??뺣━??蹂닿쿋?듬땲??

## ?щ씪?대뱶 2. ?듭떖 ?댁슜
- ?붾㈃ 臾멸뎄: 臾댁뾿???щ씪議뚮굹
- ?대젅?댁뀡: 釉붾줈洹?湲?먯꽌 ?뺣━???듭떖? ?ㅼ쓬怨?媛숈뒿?덈떎. {preview[:220]}

## ?щ씪?대뱶 3. ??以묒슂?쒓?
- ?붾㈃ 臾멸뎄: 以묒슂???댁쑀
- ?대젅?댁뀡: ???댁뒋??愿???쒖옣怨??곕━???좏깮???곹뼢??以????덇린 ?뚮Ц???먮쫫???④퍡 遊먯빞 ?⑸땲??

## ?щ씪?대뱶 4. 瑗?蹂??ъ씤??
- ?붾㈃ 臾멸뎄: 諛곌꼍, ?곹뼢, ?꾩냽 ?吏곸엫
- ?대젅?댁뀡: 諛곌꼍??臾댁뾿?몄?, ?ㅼ젣 ?곹뼢? ?대뵒源뚯??몄?, ?ㅼ쓬???대뼡 ?吏곸엫???섏삱吏 ?뺤씤?댁빞 ?⑸땲??

## ?щ씪?대뱶 5. ???諛⑸쾿
- ?붾㈃ 臾멸뎄: 吏湲?以鍮꾪븷 寃?
- ?대젅?댁뀡: ?깃툒??寃곕줎蹂대떎 ?뺤씤???뺣낫? 諛섎? 愿?먯쓣 ?④퍡 蹂닿퀬 ?먮떒?섎뒗 寃껋씠 醫뗭뒿?덈떎.

## ?щ씪?대뱶 6. 留덈Т由?
- ?붾㈃ 臾멸뎄: ?ㅼ쓬 蹂?붽? 以묒슂?⑸땲??
- ?대젅?댁뀡: ?ㅻ뒛 ?댁슜???꾩????섏뀲?ㅻ㈃ ?ㅼ쓬 ?먮쫫???④퍡 ?뺤씤??蹂댁꽭??
"""
    return {
        "tistory_post": format_tistory_post(tistory_post)[:4500],
        "thread_post": ensure_thread_length(thread_post, title, blog_post),
        "slide_script": slide_script,
        "vrew_script": slide_script,
    }


def create_local_content_package_data(item: NewsItem, article_context: str, today: str) -> dict[str, str]:
    title = _clean_news_title(item.title)
    source = item.source.strip() or "뉴스"
    link = item.link.strip()
    context = " ".join((article_context or item.summary or title).split())
    blog_post = f"""# {title}

이 뉴스는 오늘 흐름을 읽을 때 그냥 넘기기 어려운 이슈입니다.

{context}

## 왜 지금 봐야 할까

제목 하나만 보면 단순한 소식처럼 보일 수 있지만, 실제로는 정책, 시장, 생활비, 자산 판단과 연결될 수 있습니다. 그래서 지금 필요한 것은 빠른 결론보다 차분한 점검입니다.

## 확인할 것

첫째, 이 변화가 갑자기 나온 것인지 이전부터 이어진 흐름의 결과인지 확인해야 합니다.

둘째, 실제 부담이 누구에게 이동하는지 살펴야 합니다.

셋째, 다음 정책 발표와 시장 반응을 함께 봐야 합니다.

## 마무리

이번 뉴스는 하나의 신호일 수 있습니다. 후속 기사와 실제 변화를 함께 보면서 내 생활과 자산 계획에 어떤 영향을 줄지 점검해 보겠습니다.

출처: {source}
원문: {link}
확인일: {today}
"""
    derivatives = fallback_derivative_content_from_blog(blog_post, item)
    derivatives["blog_post"] = blog_post
    return derivatives

def ensure_blog_min_length(blog_post: str, item: NewsItem, article_context: str, today: str) -> str:
    blog_post = (blog_post or "").strip()
    if len(blog_post) >= 2800:
        return blog_post[:4500]

    context = " ".join((article_context or "").split())
    if not context:
        context = "?먮룞 ?섏쭛留뚯쑝濡쒕뒗 湲곗궗 蹂몃Ц??異⑸텇??媛?몄삤吏 紐삵뻽?듬땲?? ?먮Ц 留곹겕瑜??댁뼱 諛쒖뼵, ?쇱젙, 諛곌꼍 ?뺣낫瑜??뺤씤????蹂닿컯?댁빞 ?⑸땲??"

    title = item.title.strip() or "?ㅻ뒛??二쇱슂 ?댁뒪"
    source = item.source.strip() or "?댁뒪"
    link = item.link.strip()

    expansion = f"""

## 洹몃옒??吏湲???遊먯빞 ?좉퉴??
?댁뒪瑜?蹂???媛???꾩돩???쒓컙? ?쒕ぉ留?蹂닿퀬 吏?섏낀?붾뜲, 硫곗튌 ??洹??쇱씠 ?앺솢鍮꾨굹 ?ъ옄 ?щ━, ?쇱옄由? 湲곗뾽 遺꾩쐞湲곗? ?곌껐?섏뼱 ?덉뿀?ㅻ뒗 ?ъ떎???ㅻ뒭寃?源⑤떕???뚯엯?덈떎. ?대쾲 ?댁뒋??鍮꾩듂?⑸땲?? ?쒕㈃?곸쑝濡쒕뒗 ?섎굹??湲곗궗泥섎읆 蹂댁씠吏留? 洹??덉뿉???щ엺?ㅼ씠 臾댁뾿??嫄깆젙?섍퀬 臾댁뾿??湲곕??섎뒗吏 蹂댁뿬二쇰뒗 ?먮쫫???ㅼ뼱 ?덉뒿?덈떎.

?대쾲 ?뚯떇??異쒕컻?먯? {title}?낅땲?? {source}?먯꽌 ?꾪븳 ???먮쫫? ?⑥닚???ш굔 ?뚭컻??洹몄튂吏 ?딄퀬, ?욎쑝濡?愿???쒖옣怨??뺤콉, ?뚮퉬???좏깮???대뼡 ?곹뼢??以꾩? ?앷컖?섍쾶 留뚮벊?덈떎. {context[:500]}

以묒슂??寃껋? ???댁뒪瑜??⑥닚??醫뗫떎, ?섏걯?ㅻ줈 ?먮떒?섏? ?딅뒗 寃껋엯?덈떎. ?대뼡 ?댁뒋??泥섏쓬?먮뒗 ?묎쾶 蹂댁엯?덈떎. ?섏?留??쒓컙??吏?섎㈃ 鍮꾩슜, ?섏슂, ?뺤콉 諛⑺뼢, 湲곗뾽 ?꾨왂 媛숈? ?꾩떎?곸씤 臾몄젣濡??댁뼱吏????덉뒿?덈떎. 洹몃옒??吏湲??꾩슂??寃껋? ?깃툒??寃곕줎蹂대떎 李⑤텇???댁꽍?낅땲??

## ?쒕ぉ蹂대떎 諛곌꼍??癒쇱? 遊먯빞 ?⑸땲??
?댁뒪瑜?蹂???媛??癒쇱? ?뺤씤?댁빞 ??寃껋? ??吏湲????댁빞湲곌? ?섏솕?붽??낅땲?? 媛묒옄湲??깆옣??寃껋쿂??蹂댁씠???댁뒋???ㅼ젣濡쒕뒗 ?댁쟾遺???볦뿬 ???먮쫫??寃곌낵??寃쎌슦媛 留롮뒿?덈떎. ?뺤콉 蹂?? ?쒖옣 ?щ━, 湲곗뾽 ?ㅼ쟻, 援?젣 ?뺤꽭, 湲곗닠 蹂?? ?뚮퉬???됰룞 媛숈? ?붿씤??寃뱀튂硫댁꽌 ?대뒓 ?쒓컙 ?댁뒪濡??곗졇 ?섏삤???앹엯?덈떎.

?곕씪????湲곗궗瑜??쎌쓣 ?뚮룄 ?⑥닚???쒕ぉ留?蹂닿퀬 醫뗫떎, ?섏걯?ㅻ? ?먮떒?섍린蹂대떎 洹?諛곌꼍???④퍡 遊먯빞 ?⑸땲?? ?대뼡 ?댄빐愿怨꾩옄媛 ?吏곸??붿?, ?대뼡 ?쒕룄???쒖옣 議곌굔???곹뼢??以щ뒗吏, ?댁쟾 湲곗궗?ㅺ낵 鍮꾧탳?덉쓣 ???щ씪吏??먯? 臾댁뾿?몄? ?뺤씤?댁빞 ?⑸땲?? 諛곌꼍??蹂대㈃ 湲곗궗 ?섎굹媛 ?꾨땲???먮쫫??蹂댁엯?덈떎.

## ?ㅼ젣 ?곹뼢? ?대뵒源뚯? 媛덇퉴??
??踰덉㎏???곹뼢 踰붿쐞?낅땲?? 紐⑤뱺 ?댁뒪媛 紐⑤뱺 ?щ엺?먭쾶 媛숈? 臾닿쾶濡??ㅺ??ㅼ????딆뒿?덈떎. ?대뼡 ?댁뒪???뱀젙 ?낃퀎?먮쭔 ?곹뼢??二쇨퀬, ?대뼡 ?댁뒪???뚮퉬??臾쇨????ъ옄 ?щ━泥섎읆 ?곕━ ?쇱긽怨?吏곸젒 ?곌껐?⑸땲?? ???대뼡 ?댁뒪???뱀옣 ??蹂?붽? ?놁뼱 蹂댁뿬??紐??????뺤콉?대굹 ?쒖옣 媛寃⑹뿉 諛섏쁺?섍린???⑸땲??

???댁뒋??留덉갔媛吏?낅땲?? 愿???낃퀎, ?ъ옄?? ?뚮퉬?? ?뺤콉 ?대떦?먯뿉寃?媛곴컖 ?대뼡 ?섎?媛 ?덈뒗吏 ?섎닠??遊먯빞 ?⑸땲?? ?뱁엳 ?덉쓽 ?먮쫫, 鍮꾩슜 援ъ“, ?섏슂 蹂?? 洹쒖젣 媛?μ꽦, 寃쎌웳 援щ룄? ?곌껐?섎뒗 遺遺꾩씠 ?덈떎硫?洹??곹뼢? ?앷컖蹂대떎 ?ㅻ옒媛????덉뒿?덈떎.

## ?レ옄蹂대떎 諛⑺뼢??癒쇱? ?쎌뼱???⑸땲??
?댁뒪?먯꽌 ?レ옄??以묒슂?섏?留? ?レ옄 ?섎굹留뚯쑝濡??꾩껜瑜??먮떒?섎㈃ ?꾪뿕?⑸땲?? 以묒슂??寃껋? 諛⑺뼢?낅땲?? 醫뗭븘吏怨??덈뒗吏, ?섎튌吏怨??덈뒗吏, ?띾룄媛 鍮⑤씪吏?붿?, ?쒖옣??諛섏쓳???쇱떆?곸씤吏 吏?띿쟻?몄? 遊먯빞 ?⑸땲?? 援ъ껜?곸씤 ?섏튂媛 遺議깊븳 湲곗궗?쇱닔濡??붾뜑??諛⑺뼢?깆쓣 癒쇱? ?뺤씤?댁빞 ?⑸땲??

?먮Ц 湲곗궗?먯꽌 異붽?濡??뺤씤?섎㈃ 醫뗭? 寃껋? ??媛吏?낅땲?? 泥レ㎏, ?ㅼ젣 ?섏튂???쇱젙?낅땲?? ?섏㎏, 愿怨꾩옄 諛쒖뼵?낅땲?? ?뗭㎏, ?댄썑 ?꾩냽 議곗튂?낅땲?? ????媛吏媛 ?뺤씤?섎㈃ ?댁뒪??臾닿쾶媛 ?⑥뵮 ?좊챸?댁쭛?덈떎. 諛섎?濡?????媛吏媛 鍮꾩뼱 ?덈떎硫??깃툒???먮떒???쇳븯??寃껋씠 醫뗭뒿?덈떎.

## ??먭쾶 ?⑤뒗 吏덈Ц

???댁뒪瑜??쎄퀬 ?섏꽌 ??먭쾶 ?⑤뒗 吏덈Ц? 紐낇솗?⑸땲?? ??蹂?붽? ?④린?곸씤 ?吏곸엫?몄?, ?꾨땲硫?援ъ“?곸씤 蹂?붿쓽 ?쒖옉?몄? ?뺤씤?댁빞 ?⑸땲?? ?????댁뒋媛 ?뱀젙 湲곗뾽?대굹 ?낃퀎?먮쭔 ?곹뼢??二쇰뒗吏, ?꾨땲硫????볦? ?쒖옣怨??앺솢鍮? ?ъ옄 ?먮떒, ?뺤콉 諛⑺뼢源뚯? ?댁뼱吏????덈뒗吏??遊먯빞 ?⑸땲??

?낆옄???ㅼ쓬 吏덈Ц??湲곗??쇰줈 ?먮Ц???ㅼ떆 ?뺤씤?대낫硫?醫뗭뒿?덈떎.

- ???댁뒪??吏곸젒?곸씤 ?뱀궗?먮뒗 ?꾧뎄?멸??
- ?ㅼ젣 ?섏튂???쇱젙??湲곗궗 ?덉뿉 紐낇솗???쒖떆?섏뼱 ?덈뒗媛?
- 愿怨꾩옄 諛쒖뼵? ?⑥닚???꾨쭩?멸?, ?ㅽ뻾 怨꾪쉷?멸??
- ?쒖옣?대굹 ?뚮퉬?먯뿉寃??곹뼢??二쇰뒗 寃쎈줈??臾댁뾿?멸??
- ?욎쑝濡?異붽? 諛쒗몴???꾩냽 湲곗궗媛 ?섏삱 媛?μ꽦???덈뒗媛?

??吏덈Ц???듭씠 ?볦씠硫??⑥닚???댁뒪 ?뚮퉬媛 ?꾨땲???먮떒 媛?ν븳 ?뺣낫濡?諛붾앸땲??

## 釉붾줈洹?愿?먯뿉???뺣━?섎㈃

釉붾줈洹?湲?먯꽌???댁뒪瑜?洹몃?濡???린??寃껊낫???낆옄媛 ?댄빐?섍린 ?쎄쾶 ?댁꽍?섎뒗 寃껋씠 以묒슂?⑸땲?? ?쒕ぉ? 愿?ъ쓣 ?????덉뼱???섏?留? 蹂몃Ц? 怨쇱옣蹂대떎 ?뺣━媛 ?곗꽑?낅땲?? 吏湲덉쿂??湲곗궗 蹂몃Ц??異⑸텇???섏쭛?섏? ?딆? 寃쎌슦?먮뒗 ?뺤씤???댁슜怨?異붿젙 媛?ν븳 ?댁꽍??遺꾨━?댁꽌 ?곕뒗 寃껋씠 醫뗭뒿?덈떎.

?뺤씤???댁슜? 吏㏐퀬 遺꾨챸?섍쾶 ?뺣━?섍퀬, 遺議깊븳 遺遺꾩? ?먮Ц ?뺤씤 ?꾩슂?쇨퀬 ?쒖떆?섎뒗 諛⑹떇???덉쟾?⑸땲?? 洹몃옒???낆옄媛 湲???쎌쑝硫댁꽌???대뵒源뚯?媛 ?ъ떎?닿퀬 ?대뵒遺?곌? ?댁꽍?몄? 援щ텇?????덉뒿?덈떎. ?뱁엳 寃쎌젣???뺤콉 ?댁뒪???レ옄 ?섎굹媛 ?섎?瑜??ш쾶 諛붽? ???덇린 ?뚮Ц???대윴 援щ텇????以묒슂?⑸땲??

## 留덈Т由?

?대쾲 ?댁뒪???듭떖? {title}?낅땲?? ?꾩쭅 湲곗궗 蹂몃Ц ?꾩껜媛 異⑸텇???뺣낫?섏? ?딆븯?????덇린 ?뚮Ц???⑥젙?곸씤 寃곕줎???대━湲곕뒗 ?대졄?듬땲?? ?섏?留????댁뒋媛 ?섏????좏샇??媛蹂띿? ?딆뒿?덈떎. 諛곌꼍, ?곹뼢 踰붿쐞, ?꾩냽 ?吏곸엫??李⑤??濡??뺤씤?섎㈃ ?욎쑝濡쒖쓽 ?먮쫫?????뺥솗?섍쾶 ?쎌쓣 ???덉뒿?덈떎.

吏湲??④퀎?먯꽌 媛??醫뗭? ?쒕룄??鍮좊Ⅸ 寃곕줎蹂대떎 李⑤텇???뺤씤?낅땲?? ?먮Ц???댁뼱 ?섏튂? 諛쒖뼵???먭??섍퀬, ?댁뼱吏???꾩냽 湲곗궗源뚯? ?댄렣蹂몃떎硫????댁뒪媛 ?⑥닚???섎（吏쒕━ ?뚯떇?몄?, ?꾨땲硫?????蹂?붿쓽 ?쒖옉?몄? ?먮떒?섎뒗 ???꾩?????寃껋엯?덈떎.

異쒖쿂: {source}
?먮Ц: {link}
?뺤씤?? {today}
"""

    return f"{blog_post}\n{expansion}".strip()[:4500]


def fetch_article_details(url: str, timeout: int) -> tuple[str, str]:
    if not url:
        return "", ""

    request = urllib.request.Request(
        url,
        headers={
            "User-Agent": "Mozilla/5.0 MorningNewsMailer/2.0",
            "Accept-Language": "ko-KR,ko;q=0.9,en;q=0.7",
        },
    )

    try:
        with urllib.request.urlopen(request, timeout=timeout) as response:
            final_url = response.geturl()
            content_type = response.headers.get("content-type", "")
            if "text/html" not in content_type:
                return "", final_url
            raw_html = response.read(600_000).decode("utf-8", errors="replace")
    except Exception:
        return "", url

    text = extract_article_with_trafilatura(raw_html, final_url) or strip_html(raw_html)
    return text[:2500], final_url


def extract_article_with_trafilatura(raw_html: str, url: str) -> str:
    try:
        import trafilatura
    except Exception:
        return ""

    try:
        extracted = trafilatura.extract(
            raw_html,
            url=url,
            include_comments=False,
            include_tables=False,
            favor_precision=True,
        )
    except Exception:
        return ""

    return " ".join((extracted or "").split())


def fetch_news(query: str, limit: int, timeout: int = 15) -> list[NewsItem]:
    if query:
        encoded_query = urllib.parse.quote(query)
        feed_url = f"https://news.google.com/rss/search?q={encoded_query}&hl=ko&gl=KR&ceid=KR:ko"
    else:
        feed_url = DEFAULT_FEED_URL

    request = urllib.request.Request(feed_url, headers={"User-Agent": "MorningNewsMailer/2.0"})
    with urllib.request.urlopen(request, timeout=timeout) as response:
        root = ET.fromstring(response.read())

    channel = root.find("channel")
    if channel is None:
        return []

    items = []
    for item in channel.findall("item")[:limit]:
        link = item.findtext("link", default="").strip()
        description = item.findtext("description", default="")
        summary = strip_html(description)
        candidate_links = extract_links(description)
        direct_link = next((candidate for candidate in candidate_links if not is_google_news_url(candidate)), "")
        article_text, final_url = fetch_article_details(direct_link or link, timeout)
        if is_google_news_url(final_url):
            final_url = direct_link or link
        items.append(
            NewsItem(
                title=item.findtext("title", default="?쒕ぉ ?놁쓬").strip(),
                link=final_url or link,
                source=item.findtext("source", default="Google News").strip(),
                summary=summary,
                article_text=article_text,
            )
        )
    return items


def parse_slide_blocks(slide_script: str) -> list[dict[str, str]]:
    def parse_structured(value: str) -> list[dict[str, str]]:
        text = value.strip()
        if not text or text[0] not in "[{":
            return []
        try:
            parsed = json.loads(text)
        except Exception:
            try:
                parsed = ast.literal_eval(text)
            except Exception:
                return []

        if isinstance(parsed, dict):
            parsed = parsed.get("slides") or parsed.get("slide_script") or parsed.get("items") or []
        if not isinstance(parsed, list):
            return []

        structured: list[dict[str, str]] = []
        for index, item in enumerate(parsed, 1):
            if not isinstance(item, dict):
                continue
            number = item.get("slide_number") or item.get("number") or item.get("slide") or index
            title = item.get("title") or item.get("heading") or ""
            screen = item.get("screen_text") or item.get("screen") or item.get("?붾㈃ 臾멸뎄") or item.get("caption") or ""
            narration = item.get("narration") or item.get("voiceover") or item.get("?대젅?댁뀡") or item.get("script") or ""
            structured.append(
                {
                    "title": f"?щ씪?대뱶 {number}",
                    "screen": str(screen).strip() or str(title).strip(),
                    "narration": str(narration).strip(),
                }
            )
        return structured

    structured = parse_structured(slide_script)
    if structured:
        return structured

    slides: list[dict[str, str]] = []
    current: dict[str, str] | None = None
    mode = ""

    for raw_line in slide_script.splitlines():
        line = raw_line.strip()
        if not line:
            continue

        if line.startswith("##"):
            if current:
                slides.append(current)
            current = {"title": line.lstrip("# ").strip(), "screen": "", "narration": ""}
            mode = ""
            continue

        if current is None:
            continue

        if line.startswith("- ?붾㈃ 臾멸뎄:") or line.startswith("?붾㈃ 臾멸뎄:"):
            current["screen"] = line.split(":", 1)[1].strip()
            mode = ""
        elif line.startswith("- ?대젅?댁뀡:") or line.startswith("?대젅?댁뀡:"):
            current["narration"] = line.split(":", 1)[1].strip()
            mode = "narration"
        elif line == "?붾㈃ 臾멸뎄":
            mode = "screen"
        elif line == "?대젅?댁뀡":
            mode = "narration"
        elif mode == "screen":
            current["screen"] = (current["screen"] + "\n" + line).strip()
        else:
            cleaned = line.lstrip("- ").strip()
            current["narration"] = (current["narration"] + "\n" + cleaned).strip()

    if current:
        slides.append(current)

    return slides


def fallback_slide_blocks(item: NewsItem, article_context: str) -> list[dict[str, str]]:
    return [
        {
            "title": "슬라이드 1. 오늘의 뉴스",
            "screen": _clean_news_title(item.title)[:70],
            "narration": f'오늘 살펴볼 뉴스는 "{_clean_news_title(item.title)}"입니다.',
        },
        {
            "title": "슬라이드 2. 무슨 일이 있었나",
            "screen": "제목에서 읽히는 변화",
            "narration": f"이 뉴스는 단순한 사건보다 앞으로의 흐름을 보여주는 신호로 볼 수 있습니다. {article_context[:180]}",
        },
        {
            "title": "슬라이드 3. 왜 중요한가",
            "screen": "생활과 시장에 미칠 영향",
            "narration": "이 뉴스가 중요한 이유는 관련 시장뿐 아니라 우리의 생활비, 투자 판단, 정책 흐름에도 영향을 줄 수 있기 때문입니다.",
        },
        {
            "title": "슬라이드 4. 확인할 것",
            "screen": "배경, 영향, 다음 움직임",
            "narration": "변화가 시작된 배경과 실제 영향 범위, 그리고 앞으로 나올 후속 움직임을 차례대로 확인해야 합니다.",
        },
        {
            "title": "슬라이드 5. 대응 방법",
            "screen": "성급한 결론보다 점검",
            "narration": "성급히 결론을 내리기보다 원문과 반대 관점의 기사까지 함께 확인하고, 내 상황에 맞는 선택지를 준비하는 것이 좋습니다.",
        },
        {
            "title": "슬라이드 6. 마무리",
            "screen": "다음 변화가 더 중요합니다",
            "narration": "오늘의 뉴스는 끝이 아니라 시작일 수 있습니다. 이어지는 정책과 시장 반응을 계속 살펴보겠습니다.",
        },
    ]


def ensure_six_slide_blocks(slide_script: str, item: NewsItem, article_context: str) -> list[dict[str, str]]:
    parsed = parse_slide_blocks(slide_script)
    defaults = fallback_slide_blocks(item, article_context)
    blocks: list[dict[str, str]] = []

    for index in range(6):
        source = parsed[index] if index < len(parsed) else {}
        default = defaults[index]
        blocks.append(
            {
                "title": f"?щ씪?대뱶 {index + 1}",
                "screen": source.get("screen") or default["screen"],
                "narration": source.get("narration") or default["narration"],
            }
        )

    return blocks


def format_slide_script(blocks: list[dict[str, str]], source: str, link: str) -> str:
    sections = ["# 유튜브 슬라이드 대본"]
    for block in blocks:
        sections.append(
            f"""
## {block['title']}
화면 문구
{block['screen']}

내레이션
{block['narration']}
""".strip()
        )
    sections.append(f"출처: {source}\n원문: {link}")
    return "\n\n".join(sections)


def format_vrew_script(blocks: list[dict[str, str]], source: str) -> str:
    sections = ["# Vrew 대본"]
    for block in blocks:
        sections.append(
            f"""
## {block['title']}
화면 문구
{block['screen']}

내레이션
{block['narration']}
""".strip()
        )
    sections.append(f"출처: {source}")
    return "\n\n".join(sections)


def add_pptx_textbox(slide, text: str, left, top, width, height, *, size: int, bold: bool = False, color=None) -> None:
    from pptx.util import Pt

    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def create_basic_pptx(slide_script: str, output_path: Path, title: str) -> str:
    import zipfile

    def esc(value: str) -> str:
        return html.escape(value or "", quote=True)

    def emu(inches: float) -> int:
        return int(inches * 914400)

    def paragraphs(text: str, size: int, color: str = "F8FAFC", bold: bool = False) -> str:
        lines = [line.strip() for line in (text or "").splitlines() if line.strip()] or [""]
        bold_attr = ' b="1"' if bold else ""
        return "".join(
            f'<a:p><a:r><a:rPr lang="ko-KR" sz="{size * 100}"{bold_attr}>'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:rPr>'
            f"<a:t>{esc(line)}</a:t></a:r></a:p>"
            for line in lines
        )

    def shape(shape_id: int, name: str, text: str, x: float, y: float, w: float, h: float, size: int, *, bold: bool = False, fill: str = "0F172A", color: str = "F8FAFC") -> str:
        return f"""
        <p:sp>
          <p:nvSpPr><p:cNvPr id="{shape_id}" name="{esc(name)}"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>
          <p:spPr>
            <a:xfrm><a:off x="{emu(x)}" y="{emu(y)}"/><a:ext cx="{emu(w)}" cy="{emu(h)}"/></a:xfrm>
            <a:prstGeom prst="roundRect"><a:avLst/></a:prstGeom>
            <a:solidFill><a:srgbClr val="{fill}"/></a:solidFill>
            <a:ln><a:noFill/></a:ln>
          </p:spPr>
          <p:txBody><a:bodyPr wrap="square" anchor="mid" lIns="152400" rIns="152400" tIns="76200" bIns="76200"/><a:lstStyle/>{paragraphs(text, size, color, bold)}</p:txBody>
        </p:sp>
        """

    def slide_xml(number: int, slide_title: str, screen: str, narration: str) -> str:
        return f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
        <p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
          <p:cSld>
            <p:bg><p:bgPr><a:solidFill><a:srgbClr val="050A18"/></a:solidFill><a:effectLst/></p:bgPr></p:bg>
            <p:spTree>
              <p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>
              <p:grpSpPr><a:xfrm><a:off x="0" y="0"/><a:ext cx="0" cy="0"/><a:chOff x="0" y="0"/><a:chExt cx="0" cy="0"/></a:xfrm></p:grpSpPr>
              {shape(2, "Header", slide_title[:60], 0.65, 0.48, 12.0, 0.55, 17, bold=True, fill="111827")}
              {shape(3, "Screen Text", screen[:95], 0.9, 1.45, 11.4, 1.25, 33, bold=True, fill="050A18")}
              {shape(4, "Narration", "?대젅?댁뀡\\n" + narration[:520], 0.9, 3.2, 11.4, 2.55, 18, fill="0F172A")}
              {shape(5, "Footer", f"HYUNTOP NEWS | Slide {number}", 0.7, 6.88, 5.5, 0.25, 9, fill="050A18", color="94A3B8")}
            </p:spTree>
          </p:cSld>
          <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
        </p:sld>"""

    slides_data = parse_slide_blocks(slide_script)
    if not slides_data:
        slides_data = [{"title": "유튜브 슬라이드", "screen": title, "narration": slide_script[:600]}]

    all_slides = slides_data[:6]
    slide_count = len(all_slides)

    content_overrides = "\n".join(
        f'<Override PartName="/ppt/slides/slide{i}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        for i in range(1, slide_count + 1)
    )
    slide_ids = "\n".join(f'<p:sldId id="{255 + i}" r:id="rId{i}"/>' for i in range(1, slide_count + 1))
    rels = "\n".join(
        f'<Relationship Id="rId{i}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide{i}.xml"/>'
        for i in range(1, slide_count + 1)
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(output_path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr(
            "[Content_Types].xml",
            f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
            <Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
              <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
              <Default Extension="xml" ContentType="application/xml"/>
              <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
              {content_overrides}
            </Types>""",
        )
        archive.writestr(
            "_rels/.rels",
            """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
            <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
              <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
            </Relationships>""",
        )
        archive.writestr(
            "ppt/presentation.xml",
            f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
            <p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
              <p:sldIdLst>{slide_ids}</p:sldIdLst>
              <p:sldSz cx="12192000" cy="6858000" type="wide"/>
              <p:notesSz cx="6858000" cy="9144000"/>
            </p:presentation>""",
        )
        archive.writestr(
            "ppt/_rels/presentation.xml.rels",
            f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
            <Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">{rels}</Relationships>""",
        )
        for index, slide_info in enumerate(all_slides, 1):
            archive.writestr(
                f"ppt/slides/slide{index}.xml",
                slide_xml(index, slide_info.get("title", f"?щ씪?대뱶 {index}"), slide_info.get("screen", ""), slide_info.get("narration", "")),
            )

    return str(output_path)


def create_youtube_pptx(slide_script: str, output_path: Path, title: str) -> str:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Inches, Pt
    except Exception:
        return create_basic_pptx(slide_script, output_path, title)

    slides_data = parse_slide_blocks(slide_script)
    if not slides_data:
        slides_data = [{"title": "유튜브 슬라이드", "screen": title, "narration": slide_script[:600]}]

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    bg = RGBColor(5, 10, 24)
    panel = RGBColor(15, 23, 42)
    cyan = RGBColor(34, 211, 238)
    white = RGBColor(248, 250, 252)
    muted = RGBColor(148, 163, 184)

    def set_background(slide) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg

    def add_footer(slide, number: int) -> None:
        add_pptx_textbox(slide, f"HYUNTOP NEWS  |  Slide {number}", Inches(0.7), Inches(6.95), Inches(6), Inches(0.25), size=9, color=muted)

    for index, slide_info in enumerate(slides_data[:6], 1):
        slide = prs.slides.add_slide(blank_layout)
        set_background(slide)

        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.45), Inches(12.1), Inches(0.75))
        header.fill.solid()
        header.fill.fore_color.rgb = panel
        header.line.color.rgb = cyan
        header.line.width = Pt(1)
        slide_title = slide_info.get("title", f"슬라이드 {index}")
        title_size = 16 if len(slide_title) <= 38 else 13 if len(slide_title) <= 70 else 11
        add_pptx_textbox(slide, slide_title[:100], Inches(0.85), Inches(0.54), Inches(11.0), Inches(0.58), size=title_size, bold=True, color=white)

        screen_text = slide_info.get("screen") or slide_info.get("title", "")
        content_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.55), Inches(11.4), Inches(4.7))
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = panel
        content_box.line.color.rgb = RGBColor(51, 65, 85)
        screen_size = 36 if len(screen_text) <= 44 else 30 if len(screen_text) <= 80 else 24
        add_pptx_textbox(slide, screen_text[:180], Inches(1.2), Inches(2.15), Inches(10.7), Inches(3.2), size=screen_size, bold=True, color=white)
        add_footer(slide, index)

    prs.save(output_path)
    return str(output_path)


def create_gemini_blog_post_only(item: NewsItem, article_context: str, use_grounding: bool = False) -> str | None:
    global LAST_GEMINI_ERROR
    api_key = get_secret("GEMINI_API_KEY") or get_secret("GOOGLE_API_KEY")
    if not api_key:
        LAST_GEMINI_ERROR = "Gemini API 키가 없어 로컬 예비 글쓰기 엔진을 사용했습니다."
        return None

    try:
        from google import genai
        from google.genai import types
    except Exception as exc:
        LAST_GEMINI_ERROR = f"Gemini 패키지를 불러오지 못했습니다: {exc}"
        return None

    model = get_secret("GEMINI_BLOG_MODEL") or get_secret("GEMINI_MODEL_LOW_COST") or "gemini-2.5-flash"
    title = _clean_news_title(item.title)
    prompt = f"""
Write only one Korean blog article for readers in their 50s and 60s.
Return JSON only: {{"blog_post": "..."}}

Style:
- Tistory/Naver-style practical blog article, not a report.
- 2500~3500 Korean characters if possible.
- Use short paragraphs with blank lines.
- Do not say that the article text is missing.
- Do not invent exact numbers, quotes, schedules, or names that are not provided.
- If article detail is thin, expand with background, reader impact, checks, and practical interpretation.
- Do not create thread, slides, Vrew, or Tistory variants.

News title:
{title}

Source:
{item.source}

Original link:
{item.link}

Collected article/context:
{article_context}
"""

    client = genai.Client(api_key=api_key)
    config_kwargs = {
        "temperature": 0.72,
        "max_output_tokens": 4096,
        "response_mime_type": "application/json",
    }
    if use_grounding:
        config_kwargs["tools"] = [types.Tool(google_search=types.GoogleSearch())]

    blog_post = ""
    last_error = ""
    for attempt in range(3):
        retry_note = ""
        if attempt:
            retry_note = """

Previous response was invalid or too short.
Return a valid JSON object only. No markdown, no comments, no trailing comma.
Escape all line breaks inside JSON strings as \\n.
The only key must be blog_post.
"""
        try:
            response = client.models.generate_content(
                model=model,
                contents=prompt + retry_note,
                config=types.GenerateContentConfig(**config_kwargs),
            )
            data = extract_json_object(response.text)
            blog_post = normalize_generated_text_field(data.get("blog_post"))
        except Exception as exc:
            last_error = str(exc)
            LAST_GEMINI_ERROR = f"Gemini 블로그 글 응답 형식 오류로 재시도 중: {exc}"
            continue
        if len(blog_post) >= 800:
            break

    if len(blog_post) < 800:
        detail = f": {last_error}" if last_error else f": {len(blog_post)} chars"
        LAST_GEMINI_ERROR = f"Gemini blog post was too short or invalid{detail}"
        return None
    return blog_post[:4500]


def create_content_package(item: NewsItem, draft_dir_name: str, low_cost_mode: bool = False) -> ContentPackage | None:
    global LAST_GEMINI_ERROR
    today = datetime.now().strftime("%Y-%m-%d")
    has_article_text = is_useful_article_text(item.article_text, item.title)
    has_summary = is_useful_article_text(item.summary, item.title)
    article_context = ""
    if has_article_text:
        article_context = item.article_text
    elif has_summary:
        article_context = item.summary
    else:
        article_context = (
            "기사 본문 전문은 제공되지 않았습니다. 이 사실을 독자에게 직접 설명하지 말고, "
            "뉴스 제목과 출처를 글감으로 삼아 관련 배경과 독자 관점의 해설을 자연스럽게 확장하세요."
        )

    if low_cost_mode:
        blog_post = create_gemini_blog_post_only(item, article_context, use_grounding=False)
        if blog_post:
            derivatives = fallback_derivative_content_from_blog(blog_post, item)
            gemini_content = {
                "blog_post": blog_post,
                "tistory_post": derivatives["tistory_post"],
                "thread_post": derivatives["thread_post"],
                "slide_script": derivatives["slide_script"],
                "vrew_script": derivatives["vrew_script"],
            }
        else:
            gemini_content = create_local_content_package_data(item, article_context, today)
    else:
        gemini_content = create_gemini_content(item, article_context, use_grounding=not has_article_text)
        if not gemini_content:
            if not LAST_GEMINI_ERROR:
                LAST_GEMINI_ERROR = "Gemini가 품질 기준을 통과한 콘텐츠를 만들지 못했습니다."
            return None

    blog_post = gemini_content["blog_post"]
    tistory_post = gemini_content["tistory_post"]
    thread_post = gemini_content["thread_post"]
    slide_script = gemini_content["slide_script"]
    vrew_script = gemini_content["vrew_script"]

    if len(blog_post) < 1000:
        LAST_GEMINI_ERROR = f"콘텐츠가 너무 짧아 저장하지 않았습니다: {len(blog_post)}자"
        return None

    slide_blocks = ensure_six_slide_blocks(slide_script, item, article_context)
    slide_script = format_slide_script(slide_blocks, item.source, item.link)
    vrew_script = format_vrew_script(slide_blocks, item.source)

    directory_text = ""
    pptx_path = ""
    try:
        draft_dir = BASE_DIR / (draft_dir_name or "blog_drafts")
        draft_dir.mkdir(parents=True, exist_ok=True)
        safe_title = "".join(character for character in item.title if character not in '\\/:*?"<>|')
        safe_title = "-".join(safe_title.split())[:50] or "news-content"
        package_dir = draft_dir / f"{today}-{safe_title}"
        package_dir.mkdir(parents=True, exist_ok=True)
        (package_dir / "01-blog-post.md").write_text(blog_post, encoding="utf-8")
        (package_dir / "02-tistory-post.md").write_text(tistory_post, encoding="utf-8")
        (package_dir / "03-thread-post.txt").write_text(thread_post, encoding="utf-8")
        (package_dir / "04-youtube-slides.md").write_text(slide_script, encoding="utf-8")
        (package_dir / "05-vrew-script.txt").write_text(vrew_script, encoding="utf-8")
        pptx_path = create_youtube_pptx(slide_script, package_dir / "06-youtube-slides.pptx", item.title)
        directory_text = str(package_dir)
    except OSError:
        # Cloud storage may be temporary or read-only; the package still goes into the email.
        pass

    return ContentPackage(
        news_title=item.title,
        blog_post=blog_post,
        tistory_post=tistory_post,
        thread_post=thread_post,
        slide_script=slide_script,
        vrew_script=vrew_script,
        directory=directory_text,
        pptx_path=pptx_path,
    )


def build_email(
    sender: str,
    recipient: str,
    items: list[NewsItem],
    query: str,
    content_package: ContentPackage | None = None,
) -> EmailMessage:
    today = datetime.now().strftime("%Y-%m-%d")
    topic = query or "주요"
    subject = f"[아침 뉴스] {topic} 뉴스 {len(items)}개 - {today}"
    rows = "\n".join(
        f"""
        <li style="margin-bottom:18px">
          <a href="{html.escape(item.link)}" style="font-size:17px;font-weight:700;color:#155eef;text-decoration:none">
            {html.escape(item.title)}
          </a>
          <div style="margin-top:5px;color:#667085;font-size:13px">{html.escape(item.source)}</div>
        </li>
        """
        for item in items
    )
    package_html = ""
    package_plain = ""
    if content_package:
        package_html = f"""
          <hr style="border:none;border-top:1px solid #e5e7eb;margin:28px 0">
          <h2>오늘의 콘텐츠 패키지</h2>
          <p style="color:#667085">선택 뉴스: {html.escape(content_package.news_title)}</p>
          <h3>1. 블로그 글</h3>
          <pre style="white-space:pre-wrap;background:#f9fafb;border:1px solid #e5e7eb;padding:16px">{html.escape(content_package.blog_post)}</pre>
          <h3>2. 티스토리 각색 글</h3>
          <pre style="white-space:pre-wrap;background:#f9fafb;border:1px solid #e5e7eb;padding:16px">{html.escape(content_package.tistory_post)}</pre>
          <h3>3. 쓰레드 글</h3>
          <pre style="white-space:pre-wrap;background:#f9fafb;border:1px solid #e5e7eb;padding:16px">{html.escape(content_package.thread_post)}</pre>
          <h3>4. 유튜브 슬라이드 대본</h3>
          <pre style="white-space:pre-wrap;background:#f9fafb;border:1px solid #e5e7eb;padding:16px">{html.escape(content_package.slide_script)}</pre>
          <h3>5. Vrew 대본</h3>
          <pre style="white-space:pre-wrap;background:#f9fafb;border:1px solid #e5e7eb;padding:16px">{html.escape(content_package.vrew_script)}</pre>
        """
        package_plain = (
            f"\n\n[블로그 글]\n{content_package.blog_post}"
            f"\n\n[티스토리 각색 글]\n{content_package.tistory_post}"
            f"\n\n[쓰레드 글]\n{content_package.thread_post}"
            f"\n\n[유튜브 슬라이드 대본]\n{content_package.slide_script}"
            f"\n\n[Vrew 대본]\n{content_package.vrew_script}"
        )

    body = f"""
    <html lang="ko">
      <body style="background:#f5f7fa;padding:24px;font-family:Arial,'Malgun Gothic',sans-serif">
        <main style="max-width:680px;margin:auto;background:#fff;border:1px solid #e5e7eb;padding:28px">
          <h1 style="margin-top:0">?꾩묠 {html.escape(topic)} ?댁뒪</h1>
          <ol>{rows}</ol>
          {package_html}
        </main>
      </body>
    </html>
    """

    message = EmailMessage()
    message["From"] = sender
    message["To"] = recipient
    message["Subject"] = subject
    message.set_content(
        "\n".join(f"{index}. {item.title}\n{item.link}" for index, item in enumerate(items, 1)) + package_plain
    )
    message.add_alternative(body, subtype="html")
    if content_package and content_package.pptx_path:
        pptx_file = Path(content_package.pptx_path)
        if pptx_file.exists():
            message.add_attachment(
                pptx_file.read_bytes(),
                maintype="application",
                subtype="vnd.openxmlformats-officedocument.presentationml.presentation",
                filename=pptx_file.name,
            )
    return message


def send_email(message: EmailMessage, sender: str, app_password: str) -> None:
    with smtplib.SMTP_SSL("smtp.gmail.com", 465, context=ssl.create_default_context()) as smtp:
        smtp.login(sender, app_password)
        smtp.send_message(message)


def looks_broken_text(text: str) -> bool:
    if not text:
        return False
    markers = ["?댁", "?쒕", "?먮", "?좏", "?꾨", "異", "諛", "遺", "寃", "蹂", "吏", "�"]
    return any(marker in text for marker in markers)


def _clean_news_title(title: str) -> str:
    title = re.sub(r"\s+", " ", title or "").strip()
    title = re.sub(r"\s*-\s*[^-]{2,25}$", "", title).strip()
    return title or "오늘의 뉴스"


def _article_core(title: str, context: str) -> str:
    cleaned_lines = []
    for raw_line in (context or "").splitlines():
        line = raw_line.strip()
        if not line:
            continue
        line = re.sub(r"^#{1,6}\s*", "", line)
        line = re.sub(r"^[-*]\s*", "", line)
        cleaned_lines.append(line)
    text = " ".join(cleaned_lines)
    if len(text) > 120 and not looks_broken_text(text[:120]):
        return text[:900]
    return _clean_news_title(title)


def _make_reader_questions(title: str) -> list[str]:
    clean = _clean_news_title(title)
    return [
        f"{clean} 이슈가 실제 정책이나 시장 흐름으로 이어질 가능성은 어느 정도인가?",
        "직접 영향을 받는 사람은 누구이고, 부담은 어디로 이동하는가?",
        "이번 발표 이후 정부, 기업, 시장이 다음에 확인할 지점은 무엇인가?",
        "내 생활비, 자산 판단, 사업 계획과 연결되는 부분은 없는가?",
    ]


def parse_slide_blocks(slide_script: str) -> list[dict[str, str]]:
    blocks: list[dict[str, str]] = []
    current: dict[str, str] | None = None
    mode = ""

    for raw_line in (slide_script or "").splitlines():
        line = raw_line.strip().lstrip("-").strip()
        if not line:
            continue

        if line.startswith("##") and ("슬라이드" in line or "Slide" in line):
            if current:
                blocks.append(current)
            current = {"title": line.lstrip("# ").strip(), "screen": "", "narration": ""}
            mode = ""
            continue

        if current is None:
            continue

        label = line.rstrip(":：").strip()
        if label in {"화면 문구", "화면", "screen", "Screen"}:
            mode = "screen"
            continue
        if label in {"내레이션", "나레이션", "narration", "Narration"}:
            mode = "narration"
            continue
        if line.startswith(("화면 문구:", "화면:", "screen:", "Screen:")):
            mode = "screen"
            current["screen"] = (current["screen"] + "\n" + line.split(":", 1)[1].strip()).strip()
            continue
        if line.startswith(("내레이션:", "나레이션:", "narration:", "Narration:")):
            mode = "narration"
            current["narration"] = (current["narration"] + "\n" + line.split(":", 1)[1].strip()).strip()
            continue

        if mode == "screen":
            current["screen"] = (current["screen"] + "\n" + line).strip()
        elif mode == "narration":
            current["narration"] = (current["narration"] + "\n" + line).strip()
        elif not current["screen"]:
            current["screen"] = line

    if current:
        blocks.append(current)
    return blocks[:6]


def fallback_slide_blocks(item: NewsItem, article_context: str) -> list[dict[str, str]]:
    title = _clean_news_title(item.title)
    core = _article_core(item.title, article_context)
    questions = _make_reader_questions(item.title)
    return [
        {
            "title": "슬라이드 1. 오늘의 핵심 뉴스",
            "screen": title,
            "narration": f"오늘 볼 뉴스는 {title}입니다. 제목만 보고 지나치기보다 정책과 시장 흐름을 함께 읽어야 하는 이슈입니다.",
        },
        {
            "title": "슬라이드 2. 뉴스가 말하는 변화",
            "screen": core[:90],
            "narration": f"이 뉴스의 출발점은 {core[:180]}입니다. 사건 자체보다 이 흐름이 어디까지 번질 수 있는지가 중요합니다.",
        },
        {
            "title": "슬라이드 3. 왜 지금 중요한가",
            "screen": "생활비, 자산 판단, 정책 변화와 연결될 수 있습니다",
            "narration": "이런 뉴스는 처음에는 멀리 있는 이야기처럼 보이지만, 시간이 지나면 세금, 집값, 물가, 일자리, 자산 심리에 영향을 줄 수 있습니다.",
        },
        {
            "title": "슬라이드 4. 지금 확인할 것",
            "screen": "\n".join(questions[:3]),
            "narration": "성급히 결론부터 내리기보다 확인 질문을 먼저 세워야 합니다. 누가 영향을 받는지, 비용은 어디로 이동하는지, 다음 발표가 무엇인지 봐야 합니다.",
        },
        {
            "title": "슬라이드 5. 대응 관점",
            "screen": "뉴스를 소비하지 말고 선택 기준으로 바꾸기",
            "narration": "뉴스를 그대로 믿거나 흘려보내기보다 내 자산, 소비, 사업, 생활 계획에 어떤 영향을 줄 수 있는지 기준을 만들어야 합니다.",
        },
        {
            "title": "슬라이드 6. 마무리",
            "screen": "다음 변화가 어디서 시작될지 지켜보기",
            "narration": "이번 뉴스는 하나의 신호입니다. 후속 기사와 실제 정책 변화를 함께 보면서 선택지를 미리 정리해 두는 것이 중요합니다.",
        },
    ]


def ensure_six_slide_blocks(slide_script: str, item: NewsItem, article_context: str) -> list[dict[str, str]]:
    parsed = parse_slide_blocks(slide_script)
    defaults = fallback_slide_blocks(item, article_context)
    blocks: list[dict[str, str]] = []
    for index in range(6):
        source = parsed[index] if index < len(parsed) else {}
        default = defaults[index]
        title = source.get("title") or default["title"]
        if "슬라이드" not in title:
            title = f"슬라이드 {index + 1}. {title}"
        blocks.append(
            {
                "title": title,
                "screen": (source.get("screen") or default["screen"]).strip(),
                "narration": (source.get("narration") or default["narration"]).strip(),
            }
        )
    return blocks


def format_slide_script(blocks: list[dict[str, str]], source: str, link: str) -> str:
    sections = ["# 유튜브 슬라이드 대본"]
    for index, block in enumerate(blocks, 1):
        title = block["title"]
        if not title.startswith(f"슬라이드 {index}"):
            title = f"슬라이드 {index}. {title}"
        sections.append(
            f"""## {title}

화면 문구:
{block['screen']}

내레이션:
{block['narration']}"""
        )
    sections.append(f"출처: {source}\n원문: {link}")
    return "\n\n".join(sections)


def format_vrew_script(blocks: list[dict[str, str]], source: str) -> str:
    sections = ["# Vrew 대본"]
    for index, block in enumerate(blocks, 1):
        title = block["title"]
        if not title.startswith(f"슬라이드 {index}"):
            title = f"슬라이드 {index}. {title}"
        sections.append(
            f"""## {title}

화면 문구:
{block['screen']}

내레이션:
{block['narration']}"""
        )
    sections.append(f"출처: {source}")
    return "\n\n".join(sections)


def _build_tistory_post_from_blog(blog_post: str, title: str, source: str, link: str) -> str:
    clean_title = _clean_news_title(title)
    body = re.sub(r"^# .*$", "", blog_post or "", count=1, flags=re.MULTILINE).strip()
    body = re.sub(r"출처:.*", "", body).strip()
    body = re.sub(r"원문:.*", "", body).strip()
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", body) if p.strip()]
    intro = paragraphs[0] if paragraphs else f"{clean_title} 이슈를 5060 독자 관점에서 정리합니다."
    middle = "\n\n".join(paragraphs[1:5]) if len(paragraphs) > 1 else body
    questions = "\n".join(f"- {q}" for q in _make_reader_questions(title))
    tags = "#뉴스해설 #경제뉴스 #정책이슈 #부동산뉴스 #5060정보 #생활경제"
    return format_tistory_post(
        f"""# {clean_title}, 지금 확인해야 할 흐름

{intro}

## 1. 제목이 말하는 변화

이번 뉴스는 단순한 사건 전달보다 앞으로의 방향을 읽는 데 의미가 있습니다. 제목 속 표현을 보면 시장과 정책, 사람들의 선택이 어디로 움직이는지 힌트가 들어 있습니다.

{middle}

## 2. 독자가 주목해야 하는 이유

이런 뉴스는 처음에는 멀리 있는 이야기처럼 보이지만 시간이 지나면 생활비, 세금, 자산 가격, 투자 심리, 사업 환경으로 이어질 수 있습니다.

그래서 중요한 것은 지금 무엇이 바뀌고 있고, 그 변화가 어느 시점에 체감될지 확인하는 일입니다.

## 3. 바로 확인할 질문

{questions}

## 4. 오늘의 정리

지금 필요한 태도는 빠른 결론보다 차분한 확인입니다. 원문과 후속 보도를 함께 보면 이 이슈가 일시적인 소식인지, 실제 변화의 시작인지 더 분명하게 구분할 수 있습니다.

출처: {source}

원문: {link}

태그: {tags}"""
    )[:4500]


def _build_thread_post_from_blog(blog_post: str, title: str) -> str:
    clean_title = _clean_news_title(title)
    snippet = " ".join(re.sub(r"#+\s*", "", blog_post or "").split())[:180]
    text = (
        f"{clean_title} 뉴스는 제목만 보고 넘기기엔 아쉽습니다. "
        f"{snippet} "
        "핵심은 사건 자체보다 다음 변화입니다. 어떤 부담이 생기고, 어떤 정책과 시장 반응으로 이어질지 확인해야 합니다."
    )
    return ensure_thread_length(text, clean_title, blog_post)


def create_local_content_package_data(item: NewsItem, article_context: str, today: str) -> dict:
    title = _clean_news_title(item.title)
    core = _article_core(item.title, article_context)
    questions = "\n".join(f"- {q}" for q in _make_reader_questions(item.title))
    blog_post = f"""# {title}을 그냥 넘기면 안 되는 이유

뉴스를 볼 때 가장 아쉬운 순간은 제목만 보고 지나친 뒤, 며칠 뒤에야 그 일이 내 생활과 연결되어 있었다는 사실을 깨닫는 때입니다.

이번에 눈여겨볼 소식은 "{title}"입니다. 출처는 {item.source}입니다. 이 뉴스는 단순한 사건 소개가 아니라 앞으로의 정책, 시장, 생활비, 자산 판단과 연결될 수 있는 흐름을 담고 있습니다.

## 1. 제목 뒤에 숨어 있는 흐름

{core}

뉴스는 갑자기 생기는 것처럼 보이지만, 실제로는 이전부터 쌓여 온 변화가 어느 순간 표면으로 올라오는 경우가 많습니다. 그래서 제목만 보고 좋다, 나쁘다를 판단하기보다 배경과 후속 움직임을 함께 봐야 합니다.

## 2. 왜 지금 주목해야 할까

이런 이슈는 시간이 지나면 생활비, 세금, 자산 가격, 투자 심리, 사업 환경으로 이어질 수 있습니다. 특히 정책과 시장이 함께 움직이는 뉴스라면 처음에는 작게 보여도 나중에는 체감 변화가 커질 수 있습니다.

## 3. 독자에게 영향을 줄 수 있는 부분

첫째, 비용 구조가 바뀔 수 있습니다. 세금, 규제, 공급과 수요의 변화는 결국 개인의 지출과 자산 판단에 영향을 줍니다.

둘째, 시장 심리가 흔들릴 수 있습니다. 투자자와 소비자가 같은 뉴스를 다르게 받아들이면 가격과 거래량이 먼저 움직일 수 있습니다.

셋째, 후속 정책이 더 중요해질 수 있습니다. 첫 발표보다 이후의 보완책, 시행 시기, 예외 조건이 실제 영향을 결정하는 경우가 많습니다.

## 4. 지금 바로 확인할 질문

{questions}

## 5. 마무리

이번 뉴스는 오늘 하루 소비하고 끝낼 소식이 아닐 수 있습니다. 원문과 후속 보도를 함께 보면서 변화가 내 생활과 자산 판단에 어떤 신호를 주는지 차분히 확인해 볼 필요가 있습니다.

출처: {item.source}
원문: {item.link}
확인일: {today}"""
    tistory_post = _build_tistory_post_from_blog(blog_post, title, item.source, item.link)
    thread_post = _build_thread_post_from_blog(blog_post, title)
    slide_blocks = fallback_slide_blocks(item, article_context)
    slide_script = format_slide_script(slide_blocks, item.source, item.link)
    return {
        "blog_post": blog_post[:4500],
        "tistory_post": tistory_post,
        "thread_post": thread_post,
        "slide_script": slide_script,
        "vrew_script": format_vrew_script(slide_blocks, item.source),
    }


def create_derivative_content_from_blog(
    blog_post: str,
    title: str = "",
    source: str = "직접 편집",
    link: str = "",
) -> dict[str, str]:
    title = title.strip() or extract_title_from_text(blog_post, "직접 작성한 블로그 글")
    source = source.strip() or "직접 편집"
    link = link.strip()
    data: dict[str, str] | None = None
    api_key = get_secret("GEMINI_API_KEY")

    if api_key:
        try:
            from google import genai
            from google.genai import types

            model = get_secret("GEMINI_MODEL") or "gemini-2.5-pro"
            prompt = f"""
아래 블로그 글을 바탕으로 4가지 콘텐츠를 새로 작성해 주세요.
반드시 한국어 JSON 객체 하나만 반환하세요.
키는 tistory_post, thread_post, slide_script, vrew_script만 사용하세요.

공통 원칙:
- 원문 블로그 글의 구체적인 뉴스 주제, 인물, 기관, 정책명, 시장 영향을 반영하세요.
- 확인되지 않은 숫자, 발언, 일정은 새로 만들지 마세요.
- 일반론만 반복하지 말고 제목과 본문에 있는 구체 이슈를 중심으로 쓰세요.

tistory_post:
- 블로그 글을 그대로 복사하지 말고 티스토리 업로드용으로 각색하세요.
- 2500~3500자 분량, 짧은 문단, H2/H3 소제목, 목록, 마무리, 관련 태그 5~8개를 포함하세요.
- 문단마다 빈 줄을 넣고 글을 계속 이어 쓰지 마세요.

thread_post:
- 280~330자 정도로 쓰세요.
- 뉴스의 구체 주제와 지금 봐야 하는 이유, 확인할 포인트를 담으세요.

slide_script:
- 정확히 6개 슬라이드로 쓰세요.
- "## 슬라이드 1. ..."부터 "## 슬라이드 6. ..."까지 순서대로 쓰세요.
- 각 슬라이드에는 "화면 문구:"와 "내레이션:"을 포함하세요.
- 슬라이드 1은 뉴스 제목을 한눈에 보여주는 제목 슬라이드로 만드세요.
- 모든 슬라이드에 실제 뉴스 주제가 드러나야 합니다.

vrew_script:
- slide_script와 같은 6개 장면 순서로 쓰세요.
- 말로 읽기 자연스러운 내레이션 중심으로 쓰세요.
- 내레이션에서 "주요내용", "핵심메시지"라는 단어는 쓰지 마세요.

제목: {title}
출처: {source}
원문: {link}

블로그 글:
{blog_post[:7000]}
"""
            client = genai.Client(api_key=api_key)
            response = client.models.generate_content(
                model=model,
                contents=prompt,
                config=types.GenerateContentConfig(
                    temperature=0.75,
                    max_output_tokens=8192,
                    response_mime_type="application/json",
                ),
            )
            raw_data = extract_json_object(response.text)
            data = {
                "tistory_post": format_tistory_post(str(raw_data.get("tistory_post", "")).strip())[:4500],
                "thread_post": ensure_thread_length(str(raw_data.get("thread_post", "")).strip(), title, blog_post),
                "slide_script": str(raw_data.get("slide_script", "")).strip(),
                "vrew_script": str(raw_data.get("vrew_script", "")).strip(),
            }
        except Exception:
            data = None

    item = NewsItem(title=title, link=link, source=source, summary=blog_post[:500], article_text=blog_post)
    if not data or not all(data.values()):
        data = {
            "tistory_post": _build_tistory_post_from_blog(blog_post, title, source, link),
            "thread_post": _build_thread_post_from_blog(blog_post, title),
            "slide_script": "",
            "vrew_script": "",
        }

    slide_blocks = ensure_six_slide_blocks(data.get("slide_script", ""), item, blog_post)
    data["slide_script"] = format_slide_script(slide_blocks, source, link)
    data["vrew_script"] = format_vrew_script(slide_blocks, source)
    data["thread_post"] = ensure_thread_length(data["thread_post"], title, blog_post)
    data["tistory_post"] = format_tistory_post(data["tistory_post"])[:4500]
    return data


def create_gemini_content(item: NewsItem, article_context: str, use_grounding: bool = False) -> dict | None:
    global LAST_GEMINI_ERROR
    LAST_GEMINI_ERROR = ""
    api_key = get_secret("GEMINI_API_KEY")
    if not api_key:
        LAST_GEMINI_ERROR = "GEMINI_API_KEY가 없습니다."
        return None

    try:
        from google import genai
        from google.genai import types
    except Exception as exc:
        LAST_GEMINI_ERROR = f"Gemini 패키지를 불러오지 못했습니다: {exc}"
        return None

    model = get_secret("GEMINI_MODEL") or "gemini-2.5-pro"
    title = _clean_news_title(item.title)
    prompt = f"""
당신은 5060 독자를 위한 한국어 뉴스 해설 블로그 작가입니다.
아래 뉴스 정보를 바탕으로 콘텐츠 5개를 작성하세요.
반드시 JSON 객체 하나만 반환하세요.
키는 blog_post, tistory_post, thread_post, slide_script, vrew_script만 사용하세요.

전체 원칙:
- 기사 제목, 출처, 수집 본문에 있는 구체 주제를 반드시 반영하세요.
- 기사 본문이 부족해도 "본문을 못 가져왔다", "정보가 부족하다"라고 독자에게 말하지 마세요.
- 확인되지 않은 숫자, 발언, 일정, 기업명은 만들지 마세요.
- 일반론만 반복하지 말고, 이 뉴스 제목의 쟁점과 독자 관점의 해석을 중심으로 쓰세요.
- 문단은 짧게 나누고 계속 이어 쓰지 마세요.

blog_post:
- 3000~3800자 분량.
- 티스토리/네이버 블로그처럼 실전적인 뉴스 해설 글로 쓰세요.
- 제목은 후킹형으로 시작하세요.
- 번호가 있는 소제목 4~5개를 사용하세요.
- "왜 지금 봐야 하는가", "독자에게 어떤 영향이 있는가", "지금 확인할 것", "마무리"가 드러나야 합니다.
- 기사 내용이 아닌 허구 사례는 넣지 마세요.
- 각 소제목 아래에는 최소 2개 이상의 짧은 문단을 쓰세요.

tistory_post:
- blog_post를 복사하지 말고 티스토리 업로드용으로 각색하세요.
- 2800~3600자 분량.
- SEO형 제목, H2/H3 소제목, 짧은 문단, 목록, 마무리, 태그 5~8개 포함.
- blog_post와 같은 문장 순서와 같은 표현을 반복하지 마세요.
- 블로그 글과 다른 도입부, 다른 소제목, 다른 문장 흐름으로 재작성하세요.

thread_post:
- 280~330자 정도.
- 뉴스의 구체 주제, 지금 봐야 하는 이유, 독자가 확인할 포인트를 담으세요.

slide_script:
- 정확히 6개 슬라이드.
- "## 슬라이드 1. ..."부터 "## 슬라이드 6. ..."까지 순서대로 작성.
- 각 슬라이드에 "화면 문구:"와 "내레이션:" 포함.
- 슬라이드 1은 뉴스 제목이 보이는 제목 슬라이드.
- 모든 슬라이드에 이 뉴스의 구체 내용과 주제가 들어가야 합니다.

vrew_script:
- slide_script와 같은 순서의 6개 장면.
- 말로 읽기 자연스러운 내레이션 대본으로 작성.
- "주요내용", "핵심메시지"라는 단어는 쓰지 마세요.

뉴스 제목:
{title}

출처:
{item.source}

원문 링크:
{item.link}

수집된 기사 정보:
{article_context}
"""

    client = genai.Client(api_key=api_key)

    config_kwargs = {
        "temperature": 0.78,
        "max_output_tokens": 8192,
    }
    if use_grounding:
        config_kwargs["tools"] = [types.Tool(google_search=types.GoogleSearch())]
    else:
        config_kwargs["response_mime_type"] = "application/json"

    last_quality_errors: list[str] = []
    accepted_data: dict[str, str] | None = None
    for attempt in range(3):
        retry_note = ""
        if attempt:
            retry_note = f"""

이전 응답은 품질 기준을 통과하지 못했습니다: {', '.join(last_quality_errors)}
이번에는 반드시 아래 기준을 지켜 다시 작성하세요.
- blog_post는 3000자 이상, tistory_post는 2800자 이상.
- 짧은 일반론으로 채우지 말고 뉴스 제목의 핵심 쟁점을 모든 소제목에 연결.
- blog_post와 tistory_post는 도입부, 소제목, 문장 순서를 다르게 각색.
- thread_post, slide_script, vrew_script에도 뉴스 제목의 구체 주제를 반드시 포함.
- JSON 객체 하나만 반환.
- 코드블록, 설명문, 주석, 마지막 쉼표 없이 유효한 JSON만 반환.
- 문자열 안의 줄바꿈은 반드시 \\n으로 이스케이프.
"""

        try:
            response = client.models.generate_content(
                model=model,
                contents=prompt + retry_note,
                config=types.GenerateContentConfig(**config_kwargs),
            )
            data = extract_json_object(response.text)
        except Exception as exc:
            last_quality_errors = [f"Gemini JSON 형식 오류: {exc}"]
            LAST_GEMINI_ERROR = f"Gemini 응답 형식 오류로 재시도 중: {exc}"
            continue

        blog_post = normalize_generated_text_field(data.get("blog_post"))
        tistory_post = normalize_generated_text_field(data.get("tistory_post"))
        thread_post = normalize_generated_text_field(data.get("thread_post"))
        slide_script = normalize_generated_text_field(data.get("slide_script"))
        vrew_script = normalize_generated_text_field(data.get("vrew_script"))

        if not all([blog_post, tistory_post, thread_post, slide_script, vrew_script]):
            last_quality_errors = ["Gemini 응답에서 필요한 콘텐츠 일부가 비어 있음"]
            continue

        quality_errors = validate_generated_content(
            {
                "blog_post": blog_post,
                "tistory_post": tistory_post,
                "thread_post": thread_post,
                "slide_script": slide_script,
                "vrew_script": vrew_script,
            },
            item,
        )
        if quality_errors:
            last_quality_errors = quality_errors
            continue

        accepted_data = {
            "blog_post": blog_post,
            "tistory_post": tistory_post,
            "thread_post": thread_post,
            "slide_script": slide_script,
            "vrew_script": vrew_script,
        }
        break

    if not accepted_data:
        LAST_GEMINI_ERROR = "Gemini 결과물이 품질 기준을 통과하지 못했습니다: " + ", ".join(last_quality_errors)
        return None

    blog_post = accepted_data["blog_post"]
    tistory_post = accepted_data["tistory_post"]
    thread_post = accepted_data["thread_post"]
    slide_script = accepted_data["slide_script"]
    vrew_script = accepted_data["vrew_script"]

    slide_blocks = ensure_six_slide_blocks(slide_script, item, article_context)
    return {
        "blog_post": blog_post[:4500],
        "tistory_post": format_tistory_post(tistory_post)[:4500],
        "thread_post": ensure_thread_length(thread_post, item.title, blog_post),
        "slide_script": format_slide_script(slide_blocks, item.source, item.link),
        "vrew_script": format_vrew_script(slide_blocks, item.source),
    }


def clean_narration_text(text: str) -> str:
    replacements = {
        "주요내용": "내용",
        "주요 내용": "내용",
        "핵심메시지": "핵심",
        "핵심 메시지": "핵심",
        "핵심메세지": "핵심",
        "핵심 메세지": "핵심",
    }
    cleaned = text or ""
    for old, new in replacements.items():
        cleaned = cleaned.replace(old, new)
    return cleaned


def _plain_text(text: str) -> str:
    cleaned = re.sub(r"[#*_`>\-\[\]():/.,…\"'“”‘’]", " ", text or "")
    return re.sub(r"\s+", " ", cleaned).strip()


def _title_tokens(title: str) -> list[str]:
    cleaned = re.sub(r"-\s*[^-]+$", "", _clean_news_title(title))
    tokens = re.findall(r"[가-힣A-Za-z0-9]{2,}", cleaned)
    stopwords = {
        "뉴스",
        "단독",
        "속보",
        "머니투데이",
        "연합뉴스",
        "조선일보",
        "조선비즈",
        "한국경제",
        "매일경제",
        "서울경제",
    }
    return [token for token in tokens if token not in stopwords][:10]


def _has_news_signal(text: str, title: str) -> bool:
    plain = _plain_text(text)
    tokens = _title_tokens(title)
    required = max(1, min(3, len(tokens)))
    return sum(1 for token in tokens if token in plain) >= required


def _content_similarity(left: str, right: str) -> float:
    return SequenceMatcher(None, _plain_text(left)[:4000], _plain_text(right)[:4000]).ratio()


def _generic_phrase_count(text: str) -> int:
    phrases = [
        "제목만 보고",
        "단순한 사건",
        "흐름을 읽는",
        "생활비, 세금",
        "자산 판단",
        "후속 보도",
        "차분한 확인",
        "빠른 결론보다",
        "하나의 신호",
    ]
    return sum(1 for phrase in phrases if phrase in (text or ""))


def validate_generated_content(data: dict, item: NewsItem) -> list[str]:
    title = item.title
    blog_post = str(data.get("blog_post", "")).strip()
    tistory_post = str(data.get("tistory_post", "")).strip()
    thread_post = str(data.get("thread_post", "")).strip()
    slide_script = str(data.get("slide_script", "")).strip()
    vrew_script = str(data.get("vrew_script", "")).strip()
    errors: list[str] = []

    if len(blog_post) < 2800:
        errors.append(f"블로그 글이 너무 짧음({len(blog_post)}자)")
    if len(tistory_post) < 2400:
        errors.append(f"티스토리 글이 너무 짧음({len(tistory_post)}자)")
    if _content_similarity(blog_post, tistory_post) > 0.78:
        errors.append("블로그 글과 티스토리 글이 너무 비슷함")
    if not _has_news_signal(blog_post, title):
        errors.append("블로그 글에 뉴스 고유 내용이 부족함")
    if not _has_news_signal(tistory_post, title):
        errors.append("티스토리 글에 뉴스 고유 내용이 부족함")
    if not _has_news_signal(thread_post, title):
        errors.append("쓰레드 글에 뉴스 고유 내용이 부족함")
    if not _has_news_signal(slide_script, title):
        errors.append("유튜브 대본에 뉴스 고유 내용이 부족함")
    if not _has_news_signal(vrew_script, title):
        errors.append("Vrew 대본에 뉴스 고유 내용이 부족함")
    if _generic_phrase_count(blog_post) >= 4 or _generic_phrase_count(tistory_post) >= 4:
        errors.append("일반론 문장이 너무 많음")
    if "??" in blog_post + tistory_post + thread_post + slide_script + vrew_script:
        errors.append("깨진 글자가 포함됨")

    return errors


def ensure_six_slide_blocks(slide_script: str, item: NewsItem, article_context: str) -> list[dict[str, str]]:
    parsed = parse_slide_blocks(slide_script)
    defaults = fallback_slide_blocks(item, article_context)
    blocks: list[dict[str, str]] = []
    for index in range(6):
        source = parsed[index] if index < len(parsed) else {}
        default = defaults[index]
        title = source.get("title") or default["title"]
        if "슬라이드" not in title:
            title = f"슬라이드 {index + 1}. {title}"
        blocks.append(
            {
                "title": title,
                "screen": (source.get("screen") or default["screen"]).strip(),
                "narration": clean_narration_text((source.get("narration") or default["narration"]).strip()),
            }
        )
    return blocks


def format_slide_script(blocks: list[dict[str, str]], source: str, link: str) -> str:
    sections = ["# 유튜브 슬라이드 대본"]
    for index, block in enumerate(blocks, 1):
        title = block["title"]
        screen = block["screen"]
        if index == 1 and screen:
            title = f"슬라이드 1. {screen[:70]}"
        elif not title.startswith(f"슬라이드 {index}"):
            title = f"슬라이드 {index}. {title}"
        sections.append(
            f"""## {title}

화면 문구:
{screen}

내레이션:
{clean_narration_text(block['narration'])}"""
        )
    sections.append(f"출처: {source}\n원문: {link}")
    return "\n\n".join(sections)


def format_vrew_script(blocks: list[dict[str, str]], source: str) -> str:
    sections = ["# Vrew 대본"]
    for index, block in enumerate(blocks, 1):
        title = block["title"]
        if index == 1 and block.get("screen"):
            title = f"슬라이드 1. {block['screen'][:70]}"
        elif not title.startswith(f"슬라이드 {index}"):
            title = f"슬라이드 {index}. {title}"
        sections.append(
            f"""## {title}

화면 문구:
{block['screen']}

내레이션:
{clean_narration_text(block['narration'])}"""
        )
    sections.append(f"출처: {source}")
    return "\n\n".join(sections)


def run_mailer(settings: dict | None = None, dry_run: bool = False) -> dict:
    load_env_file()
    settings = settings or load_settings()
    query = str(settings.get("news_query", os.getenv("NEWS_QUERY", ""))).strip()
    limit = int(settings.get("news_limit", os.getenv("NEWS_LIMIT", "5")))
    timeout = int(settings.get("request_timeout_seconds", os.getenv("REQUEST_TIMEOUT_SECONDS", "15")))
    blog_enabled = bool(settings.get("blog_enabled", False))
    pick_index = int(settings.get("blog_pick_index", 1))
    low_cost_mode = as_bool(settings.get("low_cost_mode", os.getenv("LOW_COST_MODE", "false")), False)
    candidate_limit = int(settings.get("content_candidate_limit", os.getenv("CONTENT_CANDIDATE_LIMIT", "10")))
    fetch_limit = max(limit, candidate_limit) if blog_enabled and pick_index == 0 else limit
    items = fetch_news(query, fetch_limit, timeout)

    if not items:
        return {"ok": True, "sent": 0, "message": "가져온 뉴스가 없습니다."}

    email_items = items[:limit]
    content_package = None
    selected_index = None
    auto_selected = False
    content_message = ""
    if blog_enabled:
        selected_item, selected_index, auto_selected = select_content_item(items, pick_index)
        if auto_selected and article_score(selected_item) <= 0:
            content_message = "본문이 충분한 기사는 없었지만, 제목과 요약을 바탕으로 해설형 콘텐츠를 만들었습니다."
        content_package = create_content_package(
            selected_item,
            str(settings.get("blog_draft_dir", "blog_drafts")),
            low_cost_mode=low_cost_mode,
        )
        if content_package is not None and low_cost_mode:
            content_message = (
                "Gemini 비용 절약 모드가 켜져 있습니다. Gemini는 블로그 글만 작성하고, "
                "티스토리/쓰레드/유튜브/Vrew/PPTX 콘텐츠는 블로그 글을 바탕으로 로컬에서 만들었습니다."
            )
        if content_package is None:
            error_detail = LAST_GEMINI_ERROR or "상세 오류 없음"
            content_message = (
                "Gemini가 블로그 글을 만들지 못해 콘텐츠 패키지를 저장하지 않았습니다. "
                f"원인: {error_detail}"
            )

    if dry_run:
        return {
            "ok": True,
            "sent": 0,
            "items": [item.__dict__ for item in email_items],
            "content_package": content_package.__dict__ if content_package else None,
            "selected_index": selected_index,
            "auto_selected": auto_selected,
            "content_message": content_message,
        }

    sender = get_secret("GMAIL_ADDRESS")
    app_password = get_secret("GMAIL_APP_PASSWORD")
    recipient = str(settings.get("recipient_email") or get_secret("RECIPIENT_EMAIL") or sender).strip()
    if not sender or not app_password or not recipient:
        raise RuntimeError("Gmail 주소, 앱 비밀번호, 받는 이메일 설정이 필요합니다.")

    send_email(build_email(sender, recipient, email_items, query, content_package), sender, app_password)
    return {
        "ok": True,
        "sent": len(email_items),
        "recipient": recipient,
        "content_package": content_package.__dict__ if content_package else None,
        "selected_index": selected_index,
        "auto_selected": auto_selected,
        "content_message": content_message,
    }


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--dry-run", action="store_true")
    args = parser.parse_args()
    print(json.dumps(run_mailer(dry_run=args.dry_run), ensure_ascii=False))


if __name__ == "__main__":
    main()
